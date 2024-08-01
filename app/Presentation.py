from pptx import Presentation
import pptx.util
from PIL import Image


class PresentationRAG:
  def __init__(self):
    self.prs = Presentation()

  def title_slide(self, title, author):
    title_slide_layout = self.prs.slide_layouts[0]
    slide = self.prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "By\n" + author

  def table_of_contents(self, toc): # dict(str, dict) # TO FIX
    def subtree(k, i, l, tf):
      p = tf.add_paragraph()
      p.text = k
      p.level = i
      if len(list(l[k]))==0:
        return
      for j in l[k]:
          subtree(j, i+1, l[k], tf)

    toc_slide_layout = self.prs.slide_layouts[1]
    slide = self.prs.slides.add_slide(toc_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    title_shape.text = 'Table Of Contents'

    tf = shapes.placeholders[1].text_frame

    level = 1
    for sub_header in toc:
      p = tf.add_paragraph()
      p.text = sub_header
      p.level = level

      if len(list(toc[sub_header].keys()))==0:
        continue
      else:
        for item in toc[sub_header]:
          subtree(item, level, toc[sub_header], tf)

  def section_header(self, header):
    section_header_slide_layout = self.prs.slide_layouts[2]
    slide = self.prs.slides.add_slide(section_header_slide_layout)
    shapes = slide.shapes
    shapes.title.text = header

  def picture_with_caption(self, picture, caption):
    picture_slide_layout = self.prs.slide_layouts[8]
    slide = self.prs.slides.add_slide(picture_slide_layout)
    shapes = slide.shapes
    shapes.placeholders[1].insert_picture(picture)
    shapes.placeholders[2].text = caption

  def content_slide(self, header, content):
    content_slide_layout = self.prs.slide_layouts[1]
    slide = self.prs.slides.add_slide(content_slide_layout)
    shapes = slide.shapes
    shapes.title.text = header
    shapes.placeholders[1].text = content

  def save(self, path):
    self.prs.save(path)
